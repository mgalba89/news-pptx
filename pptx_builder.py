from io import BytesIO
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT = RGBColor(0x4A, 0x90, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0xAA, 0xAA, 0xAA)
SLIDE_NUM_COLOR = RGBColor(0x88, 0x88, 0x88)


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_accent_bar(slide):
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE rectangle value
        Inches(0), Inches(0),
        Inches(13.33), Inches(1.2),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def build_pptx(slides_data: list) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- Title slide ---
    slide = prs.slides.add_slide(blank)
    _set_bg(slide, DARK_BG)
    _add_textbox(
        slide, "News & Trends Overview",
        Inches(1), Inches(2.5), Inches(11.33), Inches(1.2),
        size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, date.today().strftime("%B %d, %Y"),
        Inches(1), Inches(3.9), Inches(11.33), Inches(0.6),
        size=20, color=MUTED, align=PP_ALIGN.CENTER,
    )

    # --- Content slides ---
    total = len(slides_data)
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank)
        is_last = i == total - 1

        if is_last:
            _set_bg(slide, DARK_BG)
            body_color = RGBColor(0xCC, 0xCC, 0xCC)
        else:
            _set_bg(slide, LIGHT_BG)
            body_color = DARK_TEXT

        _add_accent_bar(slide)
        _add_textbox(
            slide, slide_data.get("title", ""),
            Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.9),
            size=28, bold=True, color=WHITE,
        )

        # Bullet body
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(slide_data.get("bullets", [])):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(8)
            run = p.add_run()
            run.text = f"•  {bullet}"
            run.font.size = Pt(15)
            run.font.name = "Calibri"
            run.font.color.rgb = body_color

        # Slide number (1-based, offset by 1 for the title slide)
        _add_textbox(
            slide, str(i + 2),
            Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3),
            size=10, color=SLIDE_NUM_COLOR, align=PP_ALIGN.RIGHT,
        )

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
